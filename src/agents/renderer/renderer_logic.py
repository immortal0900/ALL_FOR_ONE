from __future__ import annotations
import os
import uuid
from pathlib import Path
from typing import Dict, Any, List

from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build


from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def render_cover_slide(slide, data):
    FONT = "Noto Sans CJK KR"
    COLOR_BLACK = RGBColor(0, 0, 0)
    COLOR_BLUE = RGBColor(47, 85, 151)
    

    shapes = slide.shapes

    # 상단 협회명
    box = shapes.add_textbox(Inches(0.5), Inches(0.6), Inches(9), Inches(0.8))
    p = box.text_frame.add_paragraph()
    p.text = "■ 부동산 마케팅 협회 ■"
    p.font.bold = True
    p.font.size = Pt(20)
    p.font.name = FONT
    p.font.color.rgb = COLOR_BLACK
    p.alignment = PP_ALIGN.CENTER

    # 메인 타이틀
    title = data.get("title", "")
    if title:
        box = shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(1))
        p = box.text_frame.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.name = FONT
        p.font.color.rgb = COLOR_BLACK
        p.alignment = PP_ALIGN.CENTER

    # 핵심 요약 (골드)
    lead = data.get("lead", "")
    if lead:
        box = shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.8))
        p = box.text_frame.add_paragraph()
        p.text = lead
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.name = FONT
        p.font.color.rgb = COLOR_BLUE
        p.alignment = PP_ALIGN.CENTER

    # 하단 정보
    org = data.get("org", "")
    date = data.get("date", "")
    note = data.get("note", "※ 본 보고서는 내부 검토용입니다.")
    info = f" 주관: {org}\n 작성일: {date}\n\n{note}"
    box = shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(1.2))
    p = box.text_frame.add_paragraph()
    p.text = info
    p.font.size = Pt(14)
    p.font.name = FONT
    p.font.color.rgb = COLOR_BLACK
    p.alignment = PP_ALIGN.CENTER


def render_text_slide(slide, data):
    FONT = "Noto Sans CJK KR"
    COLOR_BLACK = RGBColor(0, 0, 0)
    COLOR_BLUE = RGBColor(47, 85, 151)

    shapes = slide.shapes

    # --- Title ---
    title = data.get("title", "")
    if title:
        box = shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.9))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_top = 0     # 내부 여백 제거
        tf.margin_bottom = 0
        p = tf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.name = FONT
        p.font.color.rgb = COLOR_BLACK

    # --- Lead (Gold) ---
    lead = data.get("lead", "")
    if lead:
        box = shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(1))
        p = box.text_frame.add_paragraph()
        p.text = lead
        p.font.bold = True
        p.font.size = Pt(18)
        p.font.name = FONT
        p.font.color.rgb = COLOR_BLUE

    # --- Groups ---
    groups = data.get("groups", [])
    col_w, row_h = 4.5, 2.6
    margin_x, margin_y = 0.5, 2.0

    for i, group in enumerate(groups[:4]):
        col, row = i % 2, i // 2
        x = margin_x + col * col_w
        
        # 첫 번째 줄(row == 0)만 위로 0.4인치 당겨서 리드와 간격 좁히기
        if row == 0:
            y = margin_y - 0.55
        else:
            y = (margin_y - 0.2) + row * row_h


        label = group.get("label", "")
        insight = group.get("insight", "")
        details = group.get("details", {})

        # Main box
        box = shapes.add_textbox(Inches(x), Inches(y), Inches(col_w - 0.3), Inches(row_h - 0.2))
        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_bottom = 0
        tf.margin_left = 0
        tf.margin_right = 0
        tf.margin_top = 0

        # Label
        if label:
            p = tf.add_paragraph()
            p.text = label
            p.font.bold = True
            p.font.size = Pt(14)
            p.font.name = FONT
            p.font.color.rgb = COLOR_BLACK

        # Insight (Gold)
        if insight:
            p = tf.add_paragraph()
            p.text = insight
            p.font.size = Pt(12)
            p.font.name = FONT
            p.font.color.rgb = COLOR_BLUE

        # Details Section
        if details:
            # case 1) 최신 포맷: ["문장1", "문장2", ...]
            if isinstance(details, list):
                # 최대 4개 항목으로 제한
                limited_details = details[:4]

                p = tf.add_paragraph()
                p.text = "• 세부:"
                p.font.bold = True
                p.font.size = Pt(11)
                p.font.name = FONT
                for d in limited_details:
                    # 각 항목이 너무 길면 잘라내기 (80자 제한)
                    truncated = d[:80] + "..." if len(d) > 80 else d
                    p = tf.add_paragraph()
                    p.text = f"   - {truncated}"
                    p.font.size = Pt(9)  # 10 -> 9로 축소
                    p.font.name = FONT

            # case 2) 예전 포맷: {"data": [...], "analysis": [...], "method": "..."}
            elif isinstance(details, dict):
                data_list = details.get("data", [])
                analysis_list = details.get("analysis", [])
                method = details.get("method")

                if data_list:
                    # 최대 3개 항목으로 제한
                    limited_data = data_list[:3]

                    p = tf.add_paragraph()
                    p.text = "• 데이터:"
                    p.font.bold = True
                    p.font.size = Pt(11)
                    p.font.name = FONT
                    for d in limited_data:
                        # 각 항목이 너무 길면 잘라내기 (70자 제한)
                        truncated = d[:70] + "..." if len(d) > 70 else d
                        p = tf.add_paragraph()
                        p.text = f"   - {truncated}"
                        p.font.size = Pt(9)  # 10 -> 9로 축소
                        p.font.name = FONT

                if analysis_list:
                    # 최대 3개 항목으로 제한
                    limited_analysis = analysis_list[:3]

                    p = tf.add_paragraph()
                    p.text = "• 분석:"
                    p.font.bold = True
                    p.font.size = Pt(11)
                    p.font.name = FONT
                    for a in limited_analysis:
                        # 각 항목이 너무 길면 잘라내기 (70자 제한)
                        truncated = a[:70] + "..." if len(a) > 70 else a
                        p = tf.add_paragraph()
                        p.text = f"   - {truncated}"
                        p.font.size = Pt(9)  # 10 -> 9로 축소
                        p.font.name = FONT

from pptx.enum.text import MSO_ANCHOR
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

def render_table_slide(slide, data):
    """일반 데이터 표 렌더링 함수 (layout_type: 4)"""
    FONT = "Noto Sans CJK KR"
    COLOR_BLACK = RGBColor(0, 0, 0)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_HEADER_BG = RGBColor(47, 85, 151)

    title = data.get("title", "")
    lead = data.get("lead", "")
    table_data = data.get("table", {})

    # Title
    if title:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
        tf = box.text_frame
        p = tf.add_paragraph()
        p.text = title
        p.font.bold = True
        p.font.size = Pt(24)
        p.font.name = FONT
        p.alignment = PP_ALIGN.CENTER

    # Lead (요약문)
    start_y = 1.0
    if lead:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.9), Inches(9), Inches(0.5))
        p = box.text_frame.add_paragraph()
        p.text = lead
        p.font.bold = True
        p.font.size = Pt(14)
        p.font.name = FONT
        p.font.color.rgb = RGBColor(47, 85, 151)
        start_y = 1.4

    # Table 렌더링
    headers = table_data.get("headers", [])
    rows_data = table_data.get("rows", [])

    if headers and rows_data:
        num_rows = len(rows_data) + 1  # 헤더 포함
        num_cols = len(headers)

        # 행 수에 따라 동적으로 높이 계산
        # 기본: 헤더 0.4인치 + 각 행 0.35인치
        row_height = 0.35
        header_height = 0.4
        total_height = header_height + (len(rows_data) * row_height)
        # 최대 높이 제한 (슬라이드를 벗어나지 않도록)
        total_height = min(total_height, 5.0)

        table_shape = slide.shapes.add_table(
            num_rows, num_cols,
            Inches(0.5), Inches(start_y),
            Inches(9), Inches(total_height)
        ).table

        # 헤더 스타일링
        for col_idx, header in enumerate(headers):
            cell = table_shape.rows[0].cells[col_idx]
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = COLOR_HEADER_BG

            for p in cell.text_frame.paragraphs:
                p.font.bold = True
                p.font.size = Pt(11)
                p.font.name = FONT
                p.font.color.rgb = COLOR_WHITE
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        # 데이터 행 스타일링
        for row_idx, row_data in enumerate(rows_data, start=1):
            for col_idx, cell_data in enumerate(row_data):
                cell = table_shape.rows[row_idx].cells[col_idx]
                cell.text = str(cell_data)

                # 셀 여백 조정
                cell.text_frame.word_wrap = True
                cell.text_frame.margin_top = Pt(3)
                cell.text_frame.margin_bottom = Pt(3)
                cell.text_frame.margin_left = Pt(5)
                cell.text_frame.margin_right = Pt(5)

                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(9)
                    p.font.name = FONT
                    p.font.color.rgb = COLOR_BLACK
                    # 첫 번째 열은 중앙 정렬, 나머지는 왼쪽 정렬
                    p.alignment = PP_ALIGN.CENTER if col_idx == 0 else PP_ALIGN.LEFT
                    p.line_spacing = 1.0

                cell.vertical_anchor = MSO_ANCHOR.TOP

        # 컬럼 너비 자동 조정
        total_width = Inches(9)
        col_width = total_width / num_cols
        for col_idx in range(num_cols):
            table_shape.columns[col_idx].width = col_width


def render_swot_slide(slide, data):
    FONT = "Noto Sans CJK KR"
    COLOR_BLACK = RGBColor(0, 0, 0)
    COLOR_WHITE = RGBColor(255, 255, 255)


    title = data.get("title", "SWOT 분석")
    groups = data.get("groups", [])

    # Title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = box.text_frame
    p = tf.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(28)  # 32 -> 28로 축소
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER

    # Table (4 rows + header) - 높이를 동적으로 조정
    rows = len(groups) + 1
    cols = 2

    # 각 그룹의 내용 길이에 따라 필요한 높이 계산
    header_height = 0.4
    total_content_height = 0
    for group in groups:
        content_lines = group.get("details", [])
        # 각 항목당 약 0.25인치, 최소 0.6인치
        row_height = max(0.6, len(content_lines[:3]) * 0.28)
        total_content_height += row_height

    # 전체 테이블 높이 (헤더 + 내용, 최대 5.2인치로 제한)
    table_height = min(5.2, header_height + total_content_height)

    table_shape = slide.shapes.add_table(
        rows, cols,
        Inches(0.5), Inches(1.1),
        Inches(9), Inches(table_height)
    ).table

    # 행 높이를 내용에 맞게 조정
    table_shape.rows[0].height = Inches(header_height)
    for i, group in enumerate(groups, start=1):
        content_lines = group.get("details", [])
        row_height = max(0.6, len(content_lines[:3]) * 0.28)
        table_shape.rows[i].height = Inches(min(row_height, 1.5))

    # Header
    hdr_cells = table_shape.rows[0].cells
    hdr_cells[0].text = "구분"
    hdr_cells[1].text = "주요 내용"
    for cell in hdr_cells:
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(13)  # 14 -> 13으로 축소
            p.font.name = FONT
            p.font.color.rgb = COLOR_WHITE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Rows
    for i, group in enumerate(groups, start=1):
        label_cell, content_cell = table_shape.rows[i].cells
        label_cell.text = group.get("label", "")
        content_lines = group.get("details", [])

        # 내용이 너무 길면 잘라내기 (3개 항목으로 제한)
        if len(content_lines) > 3:
            content_lines = content_lines[:3]

        content_cell.text = "\n".join([f"- {line}" for line in content_lines])

        # 셀 여백 조정
        content_cell.text_frame.word_wrap = True
        content_cell.text_frame.margin_top = Pt(5)
        content_cell.text_frame.margin_bottom = Pt(5)
        content_cell.text_frame.margin_left = Pt(8)
        content_cell.text_frame.margin_right = Pt(8)

        for p in label_cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(12)  # 13 -> 12로 축소
            p.font.name = FONT
            p.font.color.rgb = COLOR_BLACK
            p.alignment = PP_ALIGN.CENTER

        for p in content_cell.text_frame.paragraphs:
            p.font.size = Pt(10)  # 11 -> 10으로 축소
            p.font.name = FONT
            p.font.color.rgb = COLOR_BLACK
            p.alignment = PP_ALIGN.LEFT
            p.line_spacing = 1.1  # 줄 간격 조정

    # 테이블 스타일 간격 조정
    table_shape.columns[0].width = Inches(1.8)  # 2.0 -> 1.8로 축소
    table_shape.columns[1].width = Inches(7.2)  # 7.0 -> 7.2로 확대


from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os
from utils.util import get_project_root
import uuid
def render_sliceplan_local(slice_plan: dict, output_path:str | None = None):
    if output_path is None:
        is_docker = os.path.exists("/.dockerenv") or os.getenv("DOCKER_ENV") == "true"
        if is_docker:
            output_path = str(Path("/tmp") / f"{uuid.uuid4().hex}.pptx")
        else:
            output_path = str(
                get_project_root()
                / "src" / "agents" / "renderer"
                / "temp" /f"{uuid.uuid4().hex}.pptx"
            )

    # 디렉토리가 없으면 생성
    output_dir = Path(output_path).parent
    try:
        output_dir.mkdir(parents=True, exist_ok=True)
    except OSError as e:
        print(f"경고: 디렉토리 생성 실패: {e}")
        output_path = str(Path("/tmp") / f"{uuid.uuid4().hex}.pptx")
        output_dir = Path(output_path).parent
        output_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    blank = prs.slide_layouts[6]

    for slide_data in slice_plan.get("slides", []):
        layout_type = slide_data.get("layout_type", 2)
        slide = prs.slides.add_slide(blank)

        if layout_type == 1:
            render_cover_slide(slide, slide_data)
        elif layout_type == 2:
            render_text_slide(slide, slide_data)
        elif layout_type == 3:
            render_swot_slide(slide, slide_data)
        elif layout_type == 4:
            render_table_slide(slide, slide_data)
        else:
            print(f"[WARN] Unknown layout_type={layout_type}")

    prs.save(output_path)
    print(f"PPTX saved: {output_path}")
    return output_path
